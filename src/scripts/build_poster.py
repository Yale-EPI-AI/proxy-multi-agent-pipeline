"""Build the Yale AI Day ePoster from the Yale Single Slide Template."""
from pathlib import Path
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# -----------------------------------------------------------
# Paths
# -----------------------------------------------------------
ROOT = Path("/Users/samkouteili/yale/rose/epi/multi-agent")
TEMPLATE = Path("/Users/samkouteili/yale/rose/epi/Yale Single Slide Template.pptx")
ASSETS = ROOT / "presentations" / "poster_assets"
OUT = ROOT / "presentations" / "Kouteili.Sam.pptx"

# -----------------------------------------------------------
# Palette (matches presentation2.html)
# -----------------------------------------------------------
YALE_BLUE   = RGBColor(0x19, 0x3E, 0x77)   # matches template heading color
YALE_BLUE_LIGHT = RGBColor(0x2B, 0x5F, 0x9E)
EPI_GREEN   = RGBColor(0x2E, 0x7D, 0x32)
EPI_GREEN_BRIGHT = RGBColor(0x7C, 0xD0, 0x7F)  # for white-on-dark panels
AMBER       = RGBColor(0xBF, 0x69, 0x00)
AMBER_BRIGHT = RGBColor(0xFF, 0xC1, 0x4A)
TEXT_PRIMARY = RGBColor(0x2C, 0x2C, 0x2C)
TEXT_SECONDARY = RGBColor(0x5A, 0x5A, 0x5A)
TEXT_MUTED  = RGBColor(0x8A, 0x8A, 0x8A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_MUTED = RGBColor(0xD6, 0xDE, 0xE7)
BODY_FONT = "Calibri"

# -----------------------------------------------------------
# Copy
# -----------------------------------------------------------
TITLE_MAIN = "Agentic Environmental Data Proxy Discovery"
TITLE_SUB  = "Sam Kouteili"
TITLE_AFF  = "ROSE Lab \u00d7 EPI  \u2022  Yale University"

ABSTRACT = (
    "The Yale EPI ranks 180 countries on 58 indicators, but many suffer patchy "
    "coverage, heavy imputation, and reporting lags. We introduce a two-stage "
    "multi-agent pipeline that autonomously discovers and validates data "
    "proxies \u2014 alternative signals for hard-to-measure indicators. Across "
    "five pilot indicators it generated 48 hypotheses and validated 31 (76\u202f%)."
)

INQUIRY = (
    "The EPI aspires to rank 180 countries on environmental performance, but "
    "cannot measure what it cannot observe:\n"
    "\u2022 Waste Recovery Rate: 56 of 180 countries imputed.\n"
    "\u2022 Wastewater indicators: a single year (2015) for most countries.\n"
    "\u2022 Pesticide Risk: only two global time points (2015, 2018).\n\n"
    "When official statistics are thin, the index loses its grip on the reality "
    "it aims to describe. Validated proxies offer a way in."
)

RESEARCH_Q = (
    "Can LLM agents autonomously discover, formalize, and statistically validate "
    "proxies for hard-to-measure EPI indicators \u2014 from literature to live data "
    "to verdicts \u2014 without a human in the loop?"
)

BACKGROUND = (
    "Each candidate is formalized as h = \u03c8(c, v, r) \u2014 context, variables, "
    "relationship \u2014 following DiscoveryBench. Stage 2 fetches data through "
    "nine public sources (World Bank, WHO, UN Comtrade, NASA POWER, OpenAQ, "
    "Earth Engine, GDELT, \u2026)."
)

RESULTS_HEAD = "Results"
RESULTS_BODY = (
    "Five pilot indicators (UWD, WRR, SPI, OEB, PHL):\n"
    "\u2022  48 hypotheses generated; 41 verified\n"
    "\u2022  9 confirmed \u2502 22 partially \u2502 8 rejected \u2502 2 inconclusive\n\n"
    "Strongest validated proxies (Pearson r):\n"
    "\u2022  Basic drinking-water access \u2194 UWD: r = \u22120.814\n"
    "\u2022  Open-defecation rate \u2194 UWD: r = +0.723\n"
    "\u2022  Terrestrial protected-area % \u2194 SPI: r = +0.604\n"
    "\u2022  Fossil-fuel electricity % \u2194 WRR: r = \u22120.409"
)

CONCLUSION = (
    "A general-purpose agentic pipeline can discover and validate scientific "
    "hypotheses without a human in the loop. The recipe transfers beyond EPI "
    "to any domain asking \u201cwhat else correlates with this?\u201d"
)

REFERENCES = (
    "Block, S. et al. (2024). 2024 EPI. Yale CELP.\n"
    "Majumder, B. P. et al. (2024). DiscoveryBench. arXiv:2407.01725.\n"
    "Anthropic (2025). Claude Code SDK.\n"
    "Google DeepMind (2025). Gemini Deep Research API."
)

CONTACT = (
    "Sam Kouteili\n"
    "samkouteili@gmail.com\n"
    "github.com/SamKouteili\n"
    "ROSE \u00b7 Yale University"
)

HEADLINES = [
    ("48",   "hypotheses generated"),
    ("31",   "confirmed or partially confirmed"),
    ("r = \u22120.81", "strongest validated proxy (UWD)"),
]

FIG1_CAPTION = "Fig 1. Proxy-relationship knowledge graph"
FIG2_CAPTION = "Fig 2. Per-hypothesis dashboard (UWD)"

# -----------------------------------------------------------
# Helpers
# -----------------------------------------------------------

def set_text(tf, items, body_font=BODY_FONT, body_color=TEXT_PRIMARY):
    """Rewrite a text frame from a list of (text, opts) items.

    opts = {size_pt, bold, color, align, space_after}
    First item is the heading; remaining items are paragraphs.
    """
    # Clear existing paragraphs, keep the first one (pptx requires at least one)
    tf.clear()
    first = True
    for i, (text, opts) in enumerate(items):
        opts = opts or {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        # Multi-line paragraphs: split on '\n' into separate paragraphs
        lines = text.split("\n")
        for li, line in enumerate(lines):
            if li > 0:
                p = tf.add_paragraph()
                p.alignment = opts.get("align", PP_ALIGN.LEFT)
                if "space_after" in opts:
                    p.space_after = Pt(opts["space_after"])
            # Clear existing runs on this paragraph
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            run = p.add_run()
            run.text = line
            run.font.name = opts.get("font", body_font)
            run.font.size = Pt(opts.get("size_pt", 28))
            run.font.bold = opts.get("bold", False)
            if opts.get("color") is not None:
                run.font.color.rgb = opts["color"]
            else:
                run.font.color.rgb = body_color


def heading_block(heading, body, heading_size=36, body_size=28, heading_color=YALE_BLUE,
                  body_color=TEXT_PRIMARY):
    """Return the items list for a section with heading + body."""
    items = [(heading, {"size_pt": heading_size, "bold": True, "color": heading_color,
                        "space_after": 10})]
    for block in body.split("\n\n"):
        items.append((block.strip(), {"size_pt": body_size, "color": body_color,
                                      "space_after": 6}))
    return items


def get_shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def delete_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def replace_picture(slide, shape_name, img_path):
    sh = get_shape_by_name(slide, shape_name)
    if sh is None:
        print(f"  [warn] shape {shape_name} not found")
        return None
    left, top, width, height = sh.left, sh.top, sh.width, sh.height
    delete_shape(sh)
    pic = slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
    pic.name = shape_name + "_repl"
    return pic


def replace_chart_with_picture(slide, shape_name, img_path):
    """Replace a chart shape with a picture at identical geometry."""
    sh = get_shape_by_name(slide, shape_name)
    if sh is None:
        print(f"  [warn] chart {shape_name} not found")
        return None
    left, top, width, height = sh.left, sh.top, sh.width, sh.height
    delete_shape(sh)
    pic = slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
    pic.name = shape_name + "_repl"
    return pic


def clear_text_in_shape(slide, name, replacement=""):
    sh = get_shape_by_name(slide, name)
    if sh and sh.has_text_frame:
        sh.text_frame.clear()
        if replacement:
            p = sh.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = replacement


# -----------------------------------------------------------
# Build
# -----------------------------------------------------------
prs = Presentation(str(TEMPLATE))
slide = prs.slides[0]

# -- Title block (TextBox 56) -------------------------------
title_sh = get_shape_by_name(slide, "TextBox 56")
set_text(title_sh.text_frame, [
    (TITLE_MAIN, {"size_pt": 60, "bold": True, "color": YALE_BLUE, "space_after": 4}),
    (TITLE_SUB,  {"size_pt": 32, "color": TEXT_PRIMARY, "space_after": 2}),
    (TITLE_AFF,  {"size_pt": 28, "color": TEXT_SECONDARY}),
])

# -- Delete Picture 36 (frees space in Inquiry area) --------
pic36 = get_shape_by_name(slide, "Picture 36")
if pic36: delete_shape(pic36)

# -- Resize left-column boxes so text fits at 28pt ----------
# Column runs y=5.00 to y=22.00 (17" tall). Give each section what it needs.
ab = get_shape_by_name(slide, "TextBox 33")
ab.top = Inches(5.00); ab.height = Inches(4.60)

iq = get_shape_by_name(slide, "TextBox 34")
iq.top = Inches(9.80); iq.height = Inches(6.50)

rq = get_shape_by_name(slide, "TextBox 37")
rq.top = Inches(16.50); rq.height = Inches(4.50)

# -- Abstract (TextBox 33) ----------------------------------
set_text(ab.text_frame, heading_block("Abstract", ABSTRACT))

# -- Inquiry Justification (TextBox 34) ---------------------
set_text(iq.text_frame, heading_block("The Problem", INQUIRY))

# -- Research Question (TextBox 37) -------------------------
set_text(rq.text_frame, heading_block("Research Question", RESEARCH_Q))

# -- Background and Data (TextBox 38) -----------------------
md = get_shape_by_name(slide, "TextBox 38")
md.top = Inches(5.00); md.height = Inches(3.80)
set_text(md.text_frame, heading_block("Method & Data", BACKGROUND, body_size=28))

# -- Replace Chart 44 with pipeline diagram -----------------
replace_chart_with_picture(slide, "Chart 44", ASSETS / "pipeline_diagram.png")
pipeline_sh = get_shape_by_name(slide, "Chart 44_repl")
if pipeline_sh:
    pipeline_sh.top = Inches(8.90)
    pipeline_sh.height = Inches(4.70)

# -- Results (TextBox 45) -----------------------------------
set_text(get_shape_by_name(slide, "TextBox 45").text_frame,
         heading_block(RESULTS_HEAD, RESULTS_BODY))

# -- Simplify right-middle figures column -------------------
# Strip template's zigzag descriptions.
for dead in ("TextBox 49", "TextBox 52"):
    sh = get_shape_by_name(slide, dead)
    if sh: delete_shape(sh)

# Picture 48 (Fig 1 = KG) -> TOP-RIGHT DARK PANEL (the big one)
# Dark panel Rectangle 4 lives at (30.32, 0.28) size 10.36 x 10.89
replace_picture(slide, "Picture 48", ASSETS / "kg_large.png")
pic48 = get_shape_by_name(slide, "Picture 48_repl")
if pic48:
    pic48.left   = Inches(30.50)
    pic48.top    = Inches(0.55)
    pic48.width  = Inches(10.00)
    pic48.height = Inches(9.80)

# Fig 1 caption - WHITE text on the dark blue panel
tb53 = get_shape_by_name(slide, "TextBox 53")
tb53.left = Inches(30.50); tb53.top = Inches(10.45)
tb53.width = Inches(10.00); tb53.height = Inches(0.65)
set_text(tb53.text_frame, [
    ("Fig 1.  Knowledge graph of verified proxy relationships",
     {"size_pt": 28, "bold": True, "color": WHITE}),
])

# Picture 51 (Fig 2 = UWD dashboard) -> take old KG spot (light gray col)
replace_picture(slide, "Picture 51", ASSETS / "dashboard_UWD.png")
pic51 = get_shape_by_name(slide, "Picture 51_repl")
if pic51:
    pic51.left   = Inches(20.60)
    pic51.top    = Inches(10.30)
    pic51.width  = Inches(9.50)
    pic51.height = Inches(3.00)

tb54 = get_shape_by_name(slide, "TextBox 54")
tb54.left = Inches(20.60); tb54.top = Inches(13.35)
tb54.width = Inches(9.50); tb54.height = Inches(0.55)
set_text(tb54.text_frame, [
    ("Fig 2.  Per-hypothesis dashboard  (UWD indicator)",
     {"size_pt": 28, "bold": True, "color": TEXT_MUTED}),
])

# -- Conclusion (TextBox 55) — move below Fig 2 caption -----
cc = get_shape_by_name(slide, "TextBox 55")
cc.top = Inches(14.10); cc.height = Inches(3.50)
set_text(cc.text_frame, heading_block("Takeaway", CONCLUSION))

# -- References (TextBox 63) — on dark-blue panel (#193E77) --
set_text(get_shape_by_name(slide, "TextBox 63").text_frame,
         heading_block("References", REFERENCES,
                       heading_color=WHITE, body_color=WHITE_MUTED))

# -- Contact Us (TextBox 64) — on dark-blue panel -----------
set_text(get_shape_by_name(slide, "TextBox 64").text_frame,
         heading_block("Contact", CONTACT,
                       heading_color=WHITE, body_color=WHITE))

# -- LOGO placeholder (Rectangle 15) -> affiliation text ----
logo_sh = get_shape_by_name(slide, "Rectangle 15")
if logo_sh:
    # Keep the rectangle but replace its text
    set_text(logo_sh.text_frame, [
        ("ROSE  \u00d7  EPI", {"size_pt": 44, "bold": True, "color": YALE_BLUE,
                                "align": PP_ALIGN.CENTER}),
    ])

# -- TAP HERE kiosk button (Rectangle 2) -> delete ---------
tap_sh = get_shape_by_name(slide, "Rectangle 2")
if tap_sh: delete_shape(tap_sh)

# -- Delete "Meet the Author" cluster -----------------------
for name in ("TextBox 59", "TextBox 60", "TextBox 46",
             "Rectangle 1", "Rectangle 43"):
    sh = get_shape_by_name(slide, name)
    if sh: delete_shape(sh)

# -- "At a Glance" headline panel in old KG spot (light-gray col)
# Region: x 20.60-30.10 (9.50), y 5.10 to ~9.60 (4.50)
left = Inches(20.60)
top  = Inches(5.20)
width  = Inches(9.50)
height = Inches(4.50)

hero = slide.shapes.add_textbox(left, top, width, height)
hero.name = "HeadlineFindings"
hero_tf = hero.text_frame
hero_tf.word_wrap = True
hero_tf.margin_left = hero_tf.margin_right = Pt(10)
hero_tf.margin_top = hero_tf.margin_bottom = Pt(6)

# Heading (Yale blue on light-gray panel)
set_text(hero_tf, [
    ("At a Glance", {"size_pt": 36, "bold": True, "color": YALE_BLUE, "space_after": 12}),
])

# Number rows: moderate size, green number, dark label
for big, label in HEADLINES:
    p = hero_tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)
    r1 = p.add_run()
    r1.text = big + "    "
    r1.font.name = BODY_FONT
    r1.font.size = Pt(48)
    r1.font.bold = True
    r1.font.color.rgb = EPI_GREEN
    r2 = p.add_run()
    r2.text = label
    r2.font.name = BODY_FONT
    r2.font.size = Pt(28)
    r2.font.bold = False
    r2.font.color.rgb = TEXT_PRIMARY

# Footer line
p = hero_tf.add_paragraph()
p.alignment = PP_ALIGN.LEFT
p.space_before = Pt(8)
r = p.add_run()
r.text = "Pilot indicators:  UWD \u00b7 WRR \u00b7 SPI \u00b7 OEB \u00b7 PHL"
r.font.name = BODY_FONT
r.font.size = Pt(28)
r.font.color.rgb = TEXT_SECONDARY
r.font.italic = True

# -----------------------------------------------------------
# Save
# -----------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Wrote {OUT}  ({OUT.stat().st_size/1024:.1f} KB)")
