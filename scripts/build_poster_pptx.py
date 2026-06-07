"""Build a single-slide PPTX at 40.97" x 23.04" with the rendered poster PNG full-bleed."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = Path("/Users/samkouteili/yale/rose/epi/multi-agent")
PNG = ROOT / "presentations" / "poster_assets" / "poster.png"
OUT = ROOT / "presentations" / "Kouteili.Sam.pptx"

W_IN, H_IN = 40.97, 23.04

prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)

blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)

slide.shapes.add_picture(
    str(PNG),
    left=Emu(0),
    top=Emu(0),
    width=Inches(W_IN),
    height=Inches(H_IN),
)

prs.save(str(OUT))
print(f"wrote {OUT} ({OUT.stat().st_size:,} bytes)")
print(f"  slide dims: {W_IN}\" × {H_IN}\"")
